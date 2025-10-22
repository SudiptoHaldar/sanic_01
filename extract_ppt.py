from pptx import Presentation

prs = Presentation('ui_v1.pptx')

for i, slide in enumerate(prs.slides):
    print(f'\n=== Slide {i+1} ===')
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            if shape.text.strip():
                print(f'{shape.text}')
